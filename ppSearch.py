#Python 3.x
#Jonathan Obenland 10/24/2019
#A simply program to search a directory of power points


from pptx import Presentation
from datetime import datetime
import os

def main():
    print('')
    print('A simple power point search tool by Jonathan')
    print('')
    
    while (True):
        hitCount = 0
        keyword = input("Enter word to search for (case sensitive): ")
        path = input("Enter path of power point files: ")

        #result=[]
        os.chdir(path)
        startTime = datetime.now()
        hitList = []
        print('')
        print('Script started at',startTime)
        print('')
        print('-------------')
        print('')

        for filename in os.listdir():
            
            #startTime = datetime.now()
            result=[]
            
            if not filename.endswith('.pptx'):
                continue
            presFile = Presentation(filename)
            slides = presFile.slides
            print("Searching", filename)
            print(' ')
            for slide in presFile.slides:
                slideNum = str(slides.index(slide)+1)
                for attribute in slide.shapes:
                    if not attribute.has_text_frame:
                        continue
                    for textFrame in attribute.text_frame.paragraphs:
                        for run in textFrame.runs:
                            if keyword in run.text:
                                words = slideNum,run.text
                                result.append(words)
                                if filename not in hitList:
                                    hitList.append(filename)
                                hitCount+=1
                            else:
                                continue
            if result == []:
                print("No Hits for",keyword)
                print("Check case or try agian")
                print('')
            elif result != []:
                print(result)
                print('')
        timerCalc=datetime.now() - startTime
        print('-------------')
        print('')
        print("Proccess finished in",timerCalc,"with",hitCount,"hits in",len(hitList),"powerpoints")
        print("hits were found in",hitList)
        print('')
        print('-------------')
        print('')
if __name__ == '__main__':
    main()
